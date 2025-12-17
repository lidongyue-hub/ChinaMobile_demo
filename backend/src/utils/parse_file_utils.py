import io
import logging
import pandas as pd
from typing import Dict, Any, List
from fastapi import UploadFile
import pypdf
import docx
import pptx
import zipfile
import re

logger = logging.getLogger(__name__)

async def parse_file_content(file: UploadFile) -> Dict[str, Any]:
    """
    解析上传文件内容，返回标准化格式
    """
    filename = file.filename or "unknown"
    ext = filename.split('.')[-1].lower() if '.' in filename else ""
    content = ""
    error = None

    try:
        # 读取文件内容
        file_bytes = await file.read()
        file_obj = io.BytesIO(file_bytes)
        logger.info(f"Parsing file {filename} with extension {ext}")
        match ext:
            case 'xlsx' | 'xls' | 'csv':
                content = _parse_excel(file_obj, ext)
            case 'pdf':
                content = _parse_pdf(file_obj)
            case 'docx':
                content = _parse_word(file_obj)
            case 'doc':
                content = "[注意: .doc 是旧版 Word 格式，建议转换为 .docx 后重新上传以获得更好的解析效果]"
                content += _parse_word(file_obj)
            case 'txt':
                content = file_bytes.decode('utf-8', errors='ignore')
            case 'pptx':
                content = _parse_pptx(file_obj)
            case 'ppt':
                content = "[注意: .ppt 是旧版 PowerPoint 格式，建议转换为 .pptx 后重新上传以获得更好的解析效果]"
            case _:
                logger.warning(f"Unsupported file format: {ext} for file {filename}")
                content = f"[不支持的文件格式: {ext}]"
                error = "不支持的文件格式"

    except Exception as e:
        logger.error(f"Error parsing file {filename}: {e}", exc_info=True)
        content = f"[解析失败: {str(e)}]"
        error = str(e)
    finally:
        # 关闭文件
        try:
            await file.close()
        except Exception as e:
            logger.warning(f"Error closing file {filename}: {e}")

    return {
        "name": filename,
        "type": ext,
        "content": content,
        "error": error
    }

def _parse_excel(file_obj: io.BytesIO, ext: str) -> str:
    result = []
    try:
        if ext == 'csv':
            df = pd.read_csv(file_obj)
            result.append(_dataframe_to_markdown(df, "Sheet1"))
        else:
            # xlsx or xls
            excel_file = pd.ExcelFile(file_obj)
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                result.append(_dataframe_to_markdown(df, sheet_name))
        
        content = "\n\n".join(result)
        return content if content.strip() else "[Excel 文件为空或无法读取内容]"
    except Exception as e:
        raise Exception(f"Excel解析错误: {str(e)}")

def _dataframe_to_markdown(df: pd.DataFrame, title: str) -> str:
    if df.empty:
        return ""
    
    # 限制行数，避免过长
    display_df = df.head(100)
    markdown = f"【工作表: {title}】\n"
    markdown += display_df.to_markdown(index=False)
    
    if len(df) > 100:
        markdown += f"\n... 共 {len(df)} 行数据，仅显示前 100 行"
    
    return markdown

def _parse_pdf(file_obj: io.BytesIO) -> str:
    try:
        reader = pypdf.PdfReader(file_obj)
        text_parts = []
        max_pages = min(len(reader.pages), 50)
        
        for i in range(max_pages):
            page = reader.pages[i]
            text = page.extract_text()
            if text.strip():
                text_parts.append(f"【第 {i+1} 页】\n{text}")
        
        if len(reader.pages) > max_pages:
            text_parts.append(f"\n... 共 {len(reader.pages)} 页，仅解析前 {max_pages} 页")
            
        content = "\n\n".join(text_parts)
        return content if content.strip() else "[PDF 文件为空或为扫描件（无可提取文本）]"
    except Exception as e:
        raise Exception(f"PDF解析错误: {str(e)}")

def _parse_word(file_obj: io.BytesIO) -> str:
    try:
        doc = docx.Document(file_obj)
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        
        # 也可以提取表格
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                full_text.append(" | ".join(row_text))
                
        content = "\n".join(full_text)
        return content if content.strip() else "[Word 文件为空或无法读取内容]"
    except Exception as e:
        raise Exception(f"Word解析错误: {str(e)}")

def _parse_pptx(file_obj: io.BytesIO) -> str:
    try:
        prs = pptx.Presentation(file_obj)
        text_parts = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            # 提取形状中的文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                text_parts.append(f"【幻灯片 {i+1}】\n" + "\n".join(slide_text))
                
        content = "\n\n".join(text_parts)
        return content if content.strip() else "[PPTX 文件未能提取到文本内容。该文件可能主要包含图片或图表。]"
    except Exception as e:
        raise Exception(f"PPTX解析错误: {str(e)}")
